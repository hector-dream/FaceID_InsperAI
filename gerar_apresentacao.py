"""
Gerador da apresentação da Entrega Intermediária — FaceID (Insper AI 2026.1)
===========================================================================

Cria um arquivo .pptx profissional cobrindo a rubrica:
  capa · case · benchmark · dados · arquitetura · prova de vida ·
  desafios · métricas/resultados · próximos passos

Como usar (no ambiente do projeto):
    pip install python-pptx
    python gerar_apresentacao.py

O arquivo "FaceID_Apresentacao_Intermediaria.pptx" será criado nesta pasta.
"""

import subprocess
import sys
from pathlib import Path


# ─── Bootstrap: instala python-pptx se faltar ────────────────────────────────────
def _ensure_pptx():
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass

    print("📦 python-pptx não encontrado. Tentando instalar automaticamente...")

    # Garante que o pip exista no ambiente (venvs do uv às vezes vêm sem pip).
    try:
        subprocess.check_call([sys.executable, "-m", "ensurepip", "--upgrade"])
    except Exception:
        pass

    # Tenta várias formas de instalar, em ordem.
    attempts = [
        [sys.executable, "-m", "pip", "install", "python-pptx"],
        ["uv", "pip", "install", "python-pptx"],
        ["pip", "install", "python-pptx"],
    ]
    for cmd in attempts:
        try:
            subprocess.check_call(cmd)
            import pptx  # noqa: F401
            return
        except Exception:
            continue

    sys.exit(
        "\n❌ Não consegui instalar o python-pptx automaticamente.\n"
        "   Instale manualmente e rode de novo:\n"
        "       uv pip install python-pptx\n"
        "   (ou)  python -m ensurepip --upgrade ; python -m pip install python-pptx\n"
    )


_ensure_pptx()

from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paleta (tema escuro "tech") ─────────────────────────────────────────────────
BG = RGBColor(0x0E, 0x1B, 0x2E)        # navy profundo (fundo)
CARD = RGBColor(0x1B, 0x2D, 0x45)      # navy mais claro (cartões)
CARD2 = RGBColor(0x23, 0x39, 0x56)     # navy alternativo
WHITE = RGBColor(0xF2, 0xF6, 0xFC)
MUTED = RGBColor(0xA9, 0xBB, 0xD0)     # texto secundário
ACCENT = RGBColor(0x00, 0xC9, 0xB1)    # teal/mint (destaque principal)
BLUE = RGBColor(0x3D, 0xA9, 0xFC)      # azul de apoio
RED = RGBColor(0xE5, 0x48, 0x4D)
AMBER = RGBColor(0xF5, 0xA5, 0x24)
GREEN = RGBColor(0x30, 0xC0, 0x6A)

H_FONT = "Trebuchet MS"
B_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


# ─── Helpers ─────────────────────────────────────────────────────────────────────
def slide_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def add_text(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines = lista de dicts: {text,size,color,bold,italic,font,align,space_after,line_spacing}"""
    tf = textbox(slide, l, t, w, h, anchor)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln["text"]
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get("space_after", 4))
        p.space_before = Pt(ln.get("space_before", 0))
        p.line_spacing = ln.get("line_spacing", 1.05)
        for run in p.runs:
            f = run.font
            f.size = Pt(ln.get("size", 14))
            f.bold = ln.get("bold", False)
            f.italic = ln.get("italic", False)
            f.name = ln.get("font", B_FONT)
            f.color.rgb = ln.get("color", WHITE)
    return tf


def rrect(slide, l, t, w, h, fill, line=None, line_w=1.0, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def oval(slide, l, t, w, h, fill=None, line=None, line_w=2.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def chip(slide, l, t, w, h, fill, text, tcolor=WHITE, size=12):
    sh = rrect(slide, l, t, w, h, fill, radius=0.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.name = H_FONT
        run.font.color.rgb = tcolor
    return sh


def title_block(slide, kicker, title):
    add_text(slide, 0.7, 0.5, 12, 0.4,
             [{"text": kicker, "size": 13, "color": ACCENT, "bold": True,
               "font": H_FONT}])
    add_text(slide, 0.7, 0.82, 12, 0.9,
             [{"text": title, "size": 34, "color": WHITE, "bold": True,
               "font": H_FONT}])


def footer(slide, n):
    add_text(slide, 0.7, 7.05, 8, 0.3,
             [{"text": "Insper AI · FaceID 2026.1 · Entrega Intermediária",
               "size": 9, "color": MUTED, "font": B_FONT}])
    add_text(slide, 12.0, 7.05, 0.8, 0.3,
             [{"text": str(n), "size": 9, "color": MUTED, "font": B_FONT,
               "align": PP_ALIGN.RIGHT}])


def card_with_text(slide, l, t, w, h, header, body, fill=CARD,
                   hcolor=ACCENT, header_size=15, body_size=12):
    rrect(slide, l, t, w, h, fill)
    lines = [{"text": header, "size": header_size, "color": hcolor,
              "bold": True, "font": H_FONT, "space_after": 5}]
    for b in body:
        lines.append({"text": b, "size": body_size, "color": MUTED,
                      "font": B_FONT, "space_after": 3, "line_spacing": 1.05})
    add_text(slide, l + 0.25, t + 0.22, w - 0.5, h - 0.4, lines)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

# Motivo decorativo: círculos concêntricos (vibe "scan facial") à direita
oval(s, 9.0, 1.1, 3.6, 3.6, line=CARD2, line_w=2)
oval(s, 9.6, 1.7, 2.4, 2.4, line=BLUE, line_w=2)
oval(s, 10.2, 2.3, 1.2, 1.2, line=ACCENT, line_w=2.5)
oval(s, 10.65, 2.75, 0.3, 0.3, fill=ACCENT)

chip(s, 0.9, 1.25, 3.5, 0.42, CARD2, "PROJETO FINAL · 2026.1", ACCENT, 12)
add_text(s, 0.85, 1.95, 8.2, 1.6, [
    {"text": "FaceID", "size": 72, "color": WHITE, "bold": True, "font": H_FONT},
])
add_text(s, 0.9, 3.25, 8.0, 1.2, [
    {"text": "Sistema de Reconhecimento Facial", "size": 26, "color": ACCENT,
     "bold": True, "font": H_FONT, "space_after": 2},
    {"text": "com Prova de Vida (Anti-Spoofing)", "size": 26, "color": ACCENT,
     "bold": True, "font": H_FONT},
])
add_text(s, 0.9, 4.5, 9.0, 0.6, [
    {"text": "Detecção · Embeddings · Reconhecimento · Liveness · Interface",
     "size": 14, "color": MUTED, "font": B_FONT},
])

# Equipe
rrect(s, 0.9, 5.45, 7.6, 1.05, CARD)
add_text(s, 1.15, 5.62, 7.1, 0.8, [
    {"text": "EQUIPE", "size": 11, "color": ACCENT, "bold": True,
     "font": H_FONT, "space_after": 3},
    {"text": "Hector Mathias    ·    Lara    ·    Anderson", "size": 18,
     "color": WHITE, "bold": True, "font": H_FONT},
])
add_text(s, 0.9, 6.75, 8, 0.4, [
    {"text": "Entrega Intermediária — 24/05/2026", "size": 12, "color": MUTED,
     "font": B_FONT}])


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DESCRIÇÃO DO CASE
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "O PROBLEMA", "Descrição do Case")

add_text(s, 0.7, 1.9, 11.9, 1.0, [
    {"text": "Construir, do zero e em poucos dias, um sistema de autenticação "
             "biométrica por rosto — um “Face ID” — que detecta o rosto, "
             "reconhece quem é a pessoa e confirma que é um usuário vivo "
             "(não uma foto), com uma interface de bloqueio/liberação.",
     "size": 16, "color": MUTED, "font": B_FONT, "line_spacing": 1.15}])

cy, cw, ch = 3.35, 3.83, 3.1
card_with_text(s, 0.7, cy, cw, ch, "🎯  Objetivo", [
    "Autenticar uma pessoa cadastrada em tempo real pela webcam.",
    "Liberar acesso só para o rosto correto E vivo.",
    "Simular um fluxo real de desbloqueio.",
])
card_with_text(s, 0.7 + cw + 0.2, cy, cw, ch, "🧩  Escopo", [
    "Pipeline clássico de visão: detecção → embeddings → matching.",
    "Camada de segurança: prova de vida (anti-spoofing).",
    "Protótipo funcional ponta a ponta em CPU.",
], fill=CARD2)
card_with_text(s, 0.7 + 2 * (cw + 0.2), cy, cw, ch, "⭐  Por que importa", [
    "Biometria facial está em celulares, bancos e controle de acesso.",
    "Mostra domínio de modelos pré-treinados e engenharia de pipeline.",
    "Base para virar ferramenta real da Insper AI.",
])
footer(s, 2)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BENCHMARK
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "COMO O PROBLEMA É RESOLVIDO HOJE", "Benchmark")

rows = [
    ("Apple Face ID", "Hardware + profundidade (TrueDepth/IR). Altíssima segurança, mas depende de sensor 3D dedicado."),
    ("FaceNet (Google)", "Embeddings de 128-d treinados com triplet loss. Marco do reconhecimento moderno."),
    ("ArcFace / InsightFace", "Embeddings de 512-d com margem angular aditiva. Estado da arte em precisão. ← nossa base"),
    ("dlib / face_recognition", "HOG + ResNet, 128-d. API simples, mas instalação pesada (compilação C++)."),
    ("Anti-spoofing (liveness)", "Piscada, movimento, profundidade ou CNN de textura para barrar fotos/vídeos."),
]
ty = 1.95
rh = 0.92
for i, (name, desc) in enumerate(rows):
    fill = CARD2 if name.startswith("ArcFace") else CARD
    rrect(s, 0.7, ty, 11.93, rh - 0.12, fill)
    add_text(s, 0.95, ty + 0.13, 3.4, rh - 0.3, [
        {"text": name, "size": 14, "color": ACCENT if name.startswith("ArcFace") else WHITE,
         "bold": True, "font": H_FONT}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 4.5, ty + 0.13, 8.0, rh - 0.3, [
        {"text": desc, "size": 12.5, "color": MUTED, "font": B_FONT,
         "line_spacing": 1.05}], anchor=MSO_ANCHOR.MIDDLE)
    ty += rh
footer(s, 3)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DADOS UTILIZADOS
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "DADOS E FEATURES", "Dados Utilizados")

card_with_text(s, 0.7, 1.95, 5.85, 2.25, "📸  Coleta própria (webcam)", [
    "Dataset montado por nós: data/dataset/<pessoa>/.",
    "5 a 15 fotos por pessoa em ângulos e expressões variados",
    "(augmentation natural de pose/iluminação).",
    "Recorte do rosto com margem para preservar contexto.",
])
card_with_text(s, 6.75, 1.95, 5.85, 2.25, "🧹  Pré-processamento", [
    "Detecção do rosto: MediaPipe BlazeFace.",
    "Alinhamento facial automático (InsightFace).",
    "Fallback de borda para recortes “justos”.",
    "Normalização L2 dos vetores para o matching.",
], fill=CARD2)

rrect(s, 0.7, 4.45, 11.93, 2.05, CARD)
add_text(s, 0.95, 4.65, 11.4, 1.7, [
    {"text": "🔑  Feature principal: o embedding ArcFace (512-d)", "size": 16,
     "color": ACCENT, "bold": True, "font": H_FONT, "space_after": 6},
    {"text": "Cada rosto é convertido em um vetor de 512 números. Rostos "
             "parecidos geram vetores parecidos — é isso que permite comparar "
             "identidades. Não treinamos um modelo do zero: usamos um modelo "
             "pré-treinado em milhões de rostos, e nossos dados servem para "
             "construir a base de referência (data/embeddings.pkl).",
     "size": 13.5, "color": MUTED, "font": B_FONT, "line_spacing": 1.18}])
footer(s, 4)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODELO E ARQUITETURA (pipeline)
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "MODELO E ARQUITETURA ESCOLHIDA", "Arquitetura do Pipeline")

steps = [
    ("1", "Captura", "MediaPipe\nFaceDetector"),
    ("2", "Embeddings", "InsightFace\nArcFace 512-d"),
    ("3", "Matching", "Similaridade\nde cosseno"),
    ("4", "Prova de vida", "FaceLandmarker\nEAR / piscada"),
    ("5", "Interface", "OpenCV\nbloqueio/acesso"),
]
n = len(steps)
bw, gap = 2.07, 0.30
total = n * bw + (n - 1) * gap
x = (SW - total) / 2
py = 2.2
ph = 1.85
for i, (num, head, sub) in enumerate(steps):
    rrect(s, x, py, bw, ph, CARD if i % 2 == 0 else CARD2)
    oval(s, x + bw / 2 - 0.28, py + 0.18, 0.56, 0.56, fill=ACCENT)
    add_text(s, x + bw / 2 - 0.28, py + 0.18, 0.56, 0.56, [
        {"text": num, "size": 20, "color": BG, "bold": True, "font": H_FONT,
         "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.1, py + 0.85, bw - 0.2, 0.4, [
        {"text": head, "size": 14, "color": WHITE, "bold": True, "font": H_FONT,
         "align": PP_ALIGN.CENTER}])
    add_text(s, x + 0.1, py + 1.22, bw - 0.2, 0.55, [
        {"text": sub, "size": 11, "color": MUTED, "font": B_FONT,
         "align": PP_ALIGN.CENTER, "line_spacing": 1.0}])
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(x + bw + 0.02), Inches(py + ph / 2 - 0.12),
                                Inches(gap - 0.04), Inches(0.24))
        ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
        ar.line.fill.background(); ar.shadow.inherit = False
    x += bw + gap

# Justificativas
card_with_text(s, 0.7, 4.55, 5.85, 1.95, "Por que ArcFace e não dlib?", [
    "Precisão estado-da-arte (512-d > 128-d).",
    "Modelo pré-treinado: sem custo de treino.",
    "Instala via wheels + onnxruntime (sem compilar dlib).",
    "Faz detecção e alinhamento de forma integrada.",
])
card_with_text(s, 6.75, 4.55, 5.85, 1.95, "Por que cosseno e não euclidiana?", [
    "Embeddings já são normalizados (norma 1).",
    "Cosseno vira um produto escalar — rápido e estável.",
    "Limiar interpretável (0–1); usamos 0.40.",
    "Equivalente à distância citada na rubrica (dlib ~0.6).",
], fill=CARD2)
footer(s, 5)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROVA DE VIDA / ANTI-SPOOFING
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "SEGURANÇA", "Prova de Vida (Anti-Spoofing)")

add_text(s, 0.7, 1.9, 11.9, 0.9, [
    {"text": "Um Face ID de verdade não pode ser enganado por uma foto na tela "
             "do celular. Exigimos uma piscada para provar que há uma pessoa "
             "viva — uma foto estática nunca pisca.",
     "size": 15.5, "color": MUTED, "font": B_FONT, "line_spacing": 1.15}])

card_with_text(s, 0.7, 3.0, 5.85, 2.0, "👁️  Como detectamos a piscada", [
    "Marcos faciais do MediaPipe FaceLandmarker (478 pontos).",
    "EAR (Eye Aspect Ratio) = altura / largura do olho.",
    "Olho fecha → EAR despenca; reabre → conta 1 piscada.",
])
rrect(s, 6.75, 3.0, 5.85, 2.0, CARD2)
add_text(s, 7.0, 3.25, 5.35, 1.6, [
    {"text": "Fórmula do EAR", "size": 15, "color": ACCENT, "bold": True,
     "font": H_FONT, "space_after": 8},
    {"text": "EAR = (‖p2−p6‖ + ‖p3−p5‖) / (2·‖p1−p4‖)", "size": 15,
     "color": WHITE, "bold": True, "font": "Consolas", "space_after": 8},
    {"text": "Limiar: 0.21 (abaixo = olho fechado).", "size": 12.5,
     "color": MUTED, "font": B_FONT},
])

# Fluxo de estados (chips coloridos)
add_text(s, 0.7, 5.25, 11.9, 0.35, [
    {"text": "FLUXO DE AUTENTICAÇÃO", "size": 12, "color": ACCENT, "bold": True,
     "font": H_FONT}])
chip(s, 0.7, 5.7, 3.5, 0.7, RED, "🔒  BLOQUEADO\nrosto não confirmado", WHITE, 12)
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.35), Inches(5.92),
                        Inches(0.55), Inches(0.26))
ar.fill.solid(); ar.fill.fore_color.rgb = MUTED; ar.line.fill.background(); ar.shadow.inherit = False
chip(s, 5.05, 5.7, 3.5, 0.7, AMBER, "👁️  PISQUE\nprova de vida", BG, 12)
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.7), Inches(5.92),
                        Inches(0.55), Inches(0.26))
ar.fill.solid(); ar.fill.fore_color.rgb = MUTED; ar.line.fill.background(); ar.shadow.inherit = False
chip(s, 9.4, 5.7, 3.5, 0.7, GREEN, "✅  ACESSO LIBERADO\nbem-vindo!", BG, 12)
footer(s, 6)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DESAFIOS
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "O QUE ENFRENTAMOS", "Desafios Encontrados")

challenges = [
    ("🐍  Python 3.14 vs bibliotecas",
     "onnxruntime/InsightFace sem wheels para 3.14.",
     "Resolvido migrando o ambiente para Python 3.11."),
    ("🪟  InsightFace no Windows",
     "Instalação podia exigir compilador C++.",
     "Documentado o passo a passo + uso de wheels."),
    ("✂️  Detecção em recortes justos",
     "O detector falhava em rostos sem margem.",
     "Fallback que adiciona borda/contexto à imagem."),
    ("⚡  Performance (0.7 FPS)",
     "Rodar o modelo pesado em todo quadro travava o vídeo.",
     "Reconhecimento em thread → saltou para ~58 FPS."),
]
cw, ch = 5.85, 1.95
positions = [(0.7, 1.95), (6.75, 1.95), (0.7, 4.05), (6.75, 4.05)]
for (px, py), (head, prob, sol) in zip(positions, challenges):
    rrect(s, px, py, cw, ch, CARD)
    add_text(s, px + 0.25, py + 0.2, cw - 0.5, ch - 0.4, [
        {"text": head, "size": 15, "color": WHITE, "bold": True, "font": H_FONT,
         "space_after": 6},
        {"text": "Desafio: " + prob, "size": 12, "color": MUTED, "font": B_FONT,
         "space_after": 4, "line_spacing": 1.05},
        {"text": "Solução: " + sol, "size": 12, "color": ACCENT, "font": B_FONT,
         "line_spacing": 1.05},
    ])
footer(s, 7)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MÉTRICAS E RESULTADOS
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "O QUE VAMOS MEDIR (E O QUE JÁ TEMOS)", "Métricas e Resultados")

# Stats já atingidos
stats = [
    ("0.92", "similaridade de\ncosseno (match)"),
    ("~58", "FPS em tempo\nreal (CPU)"),
    ("512", "dimensões do\nembedding"),
    ("OK", "anti-spoofing\npor piscada"),
]
sw_, gap_ = 2.85, 0.20
x = (SW - (4 * sw_ + 3 * gap_)) / 2
for val, lab in stats:
    rrect(s, x, 1.95, sw_, 1.7, CARD2)
    add_text(s, x + 0.1, 2.12, sw_ - 0.2, 0.9, [
        {"text": val, "size": 44, "color": ACCENT, "bold": True, "font": H_FONT,
         "align": PP_ALIGN.CENTER}])
    add_text(s, x + 0.1, 3.05, sw_ - 0.2, 0.5, [
        {"text": lab, "size": 12, "color": MUTED, "font": B_FONT,
         "align": PP_ALIGN.CENTER, "line_spacing": 1.0}])
    x += sw_ + gap_

card_with_text(s, 0.7, 3.95, 5.85, 2.5, "📐  Métricas-alvo (entrega final)", [
    "Acurácia de verificação (mesmo/diferente).",
    "FAR (falsa aceitação) e FRR (falsa rejeição) baixos.",
    "Limiar de cosseno calibrado (atual: 0.40).",
    "Taxa de detecção de piscada (liveness).",
    "FPS estável em tempo real.",
])
card_with_text(s, 6.75, 3.95, 5.85, 2.5, "✅  Resultado esperado", [
    "Reconhecer corretamente os cadastrados e rejeitar desconhecidos.",
    "Bloquear ataques de foto via prova de vida.",
    "Rodar fluido em CPU comum, sem GPU.",
    "Interface clara de bloqueio → liberação.",
], fill=CARD, hcolor=GREEN)
footer(s, 8)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PRÓXIMOS PASSOS / ENCERRAMENTO
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

oval(s, 9.6, 0.7, 3.4, 3.4, line=CARD2, line_w=2)
oval(s, 10.1, 1.2, 2.4, 2.4, line=BLUE, line_w=2)
oval(s, 10.6, 1.7, 1.4, 1.4, line=ACCENT, line_w=2.5)

add_text(s, 0.7, 0.95, 8.5, 0.4, [
    {"text": "DA INTERMEDIÁRIA PARA A FINAL", "size": 13, "color": ACCENT,
     "bold": True, "font": H_FONT}])
add_text(s, 0.7, 1.3, 9.0, 0.9, [
    {"text": "Próximos Passos", "size": 36, "color": WHITE, "bold": True,
     "font": H_FONT}])

nexts = [
    "Ampliar a base: mais pessoas e mais fotos por pessoa.",
    "Calibrar o limiar com métricas reais (FAR/FRR).",
    "Robustez a iluminação, óculos e variação de pose.",
    "Avaliar anti-spoofing mais forte (CNN de textura/profundidade).",
    "Escrever o documento estilo paper científico.",
]
ty = 2.5
for i, txt in enumerate(nexts):
    oval(s, 0.75, ty + 0.04, 0.28, 0.28, fill=ACCENT)
    add_text(s, 0.75, ty + 0.02, 0.28, 0.28, [
        {"text": str(i + 1), "size": 12, "color": BG, "bold": True,
         "font": H_FONT, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.25, ty, 8.0, 0.45, [
        {"text": txt, "size": 14.5, "color": MUTED, "font": B_FONT}])
    ty += 0.62

rrect(s, 0.7, 6.35, 9.0, 0.75, CARD)
add_text(s, 0.95, 6.5, 8.6, 0.5, [
    {"text": "Equipe:  Hector Mathias · Lara · Anderson      |      Insper AI 2026.1",
     "size": 13, "color": WHITE, "bold": True, "font": H_FONT}],
    anchor=MSO_ANCHOR.MIDDLE)


# ─── Salvar ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "FaceID_Apresentacao_Intermediaria.pptx"
prs.save(str(out))
print(f"\n✅ Apresentação criada: {out}")
print(f"   {len(list(prs.slides))} slides gerados.")
